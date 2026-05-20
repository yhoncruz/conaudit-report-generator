import pandas as pd
from pptx import Presentation

# Inspect Excel
excel_path = 'Reporte_Cierre_Paquete.xlsx'
df = pd.read_excel(excel_path, sheet_name='Resumen por Razón Social')
print("--- Excel Columns ---")
print(df.columns.tolist())
print(df.head(2))

# Inspect PPTX
pptx_path = 'Presentación Reunión de Seguimiento.pptx'
prs = Presentation(pptx_path)
print("\n--- PPTX Slides ---")
for i, slide in enumerate(prs.slides):
    print(f"Slide {i}:")
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(f"  - Text: {shape.text.strip()}")
        else:
            print(f"  - Shape: {shape.name}")
