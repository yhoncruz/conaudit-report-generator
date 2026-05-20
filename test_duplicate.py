from pptx import Presentation
import copy
import os

def duplicate_slide(pres, index):
    source_slide = pres.slides[index]
    # Usar un layout en blanco (usualmente el 6, o el que tenga menos cosas)
    try:
        slide_layout = pres.slide_layouts[6] 
    except:
        slide_layout = pres.slide_layouts[0]
        
    new_slide = pres.slides.add_slide(slide_layout)
    
    # Limpiar cualquier shape que venga del layout
    for shape in new_slide.shapes:
        sp = shape.element
        sp.getparent().remove(sp)
        
    # Copiar shapes del original
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(new_el)
        
    return new_slide

def main():
    pptx_path = 'Presentación Reunión de Seguimiento.pptx'
    if not os.path.exists(pptx_path):
        print("Archivo no encontrado")
        return
        
    prs = Presentation(pptx_path)
    
    print("Duplicando Slide 0...")
    new_slide = duplicate_slide(prs, 0)
    
    # Intentar cambiar el texto en el nuevo slide
    for shape in new_slide.shapes:
        if shape.has_text_frame and "Resumen Ejecutivo" in shape.text:
            if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                shape.text_frame.paragraphs[0].runs[0].text = "Resumen Ejecutivo de Gestión de\nPRUEBA DE DUPLICACIÓN"
                for i in range(1, len(shape.text_frame.paragraphs[0].runs)):
                    shape.text_frame.paragraphs[0].runs[i].text = ""
            else:
                shape.text = "Resumen Ejecutivo de Gestión de\nPRUEBA DE DUPLICACIÓN"
                
    prs.save('prueba_duplicado.pptx')
    print("Guardado como prueba_duplicado.pptx")

if __name__ == '__main__':
    main()
