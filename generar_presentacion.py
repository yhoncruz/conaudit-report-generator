import pandas as pd
import dataframe_image as dfi
from pptx import Presentation
from pptx.util import Inches
import copy
import os
import shutil

def duplicate_slide(pres, index):
    source_slide = pres.slides[index]
    # Usar un layout en blanco o el actual
    try:
        slide_layout = pres.slide_layouts[6] 
    except:
        slide_layout = pres.slide_layouts[0]
        
    new_slide = pres.slides.add_slide(slide_layout)
    
    # Limpiar cualquier shape del layout (como "Haga clic para agregar...")
    for shape in new_slide.shapes:
        sp = shape.element
        sp.getparent().remove(sp)
        
    # Copiar shapes del original al nuevo
    for shape in source_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(new_el)
        
    return new_slide

def stylize_df(df_row, columns):
    df_styled = pd.DataFrame([df_row], columns=columns)
    
    # Formateadores manuales para asegurar compatibilidad total
    format_dict = {
        'Total_Reclamado': '$',
        'Total_Aprobado': '$',
        'Total_Glosado': '$',
        'Porcentaje_Aprobado': '%',
        'Porcentaje_Glosado': '%'
    }
    
    def custom_format(x, fmt):
        if pd.isna(x):
            return ""
        
        # Para formato moneda
        if fmt == '$':
            try:
                # Tratar de manejar si viene con coma
                val_str = str(x).replace(",", ".")
                val = float(val_str)
                return "${:,.0f}".format(val).replace(",", ".")
            except:
                return str(x)
                
        # Para formato porcentaje
        elif fmt == '%':
            val_str = str(x).strip()
            if val_str.endswith('%'):
                return val_str
            # Reemplazar punto por coma según convención local
            val_str = val_str.replace(".", ",")
            return f"{val_str}%"
            
        return str(x)

    for col, fmt in format_dict.items():
        if col in df_styled.columns:
            df_styled[col] = df_styled[col].apply(lambda x: custom_format(x, fmt))

    styles = [
        dict(selector="th", props=[("background-color", "#95B3D7"),
                                   ("color", "black"),
                                   ("border", "1px solid black"),
                                   ("text-align", "center"),
                                   ("font-weight", "bold"),
                                   ("font-family", "Calibri"),
                                   ("font-size", "11pt")]),
        dict(selector="td", props=[("border", "1px solid black"),
                                   ("text-align", "right"),
                                   ("font-family", "Calibri"),
                                   ("font-size", "11pt")]),
        dict(selector="td:nth-child(1)", props=[("text-align", "left")]),
    ]
    
    return df_styled.style.hide(axis="index").set_table_styles(styles)

def main():
    excel_path = 'Reporte_Cierre_Paquete.xlsx'
    pptx_path = 'Presentación Reunión de Seguimiento.pptx'
    out_pptx_path = 'Presentacion_Actualizada_v2.pptx'
    
    print("Leyendo Excel...")
    df = pd.read_excel(excel_path, sheet_name='Resumen por Razón Social')
    
    cols_to_keep = ['Razon_Social_Reclamante', 'Cantidad', 'Total_Reclamado', 'Total_Aprobado', 'Total_Glosado', 'Porcentaje_Aprobado', 'Porcentaje_Glosado']
    existing_cols = [c for c in cols_to_keep if c in df.columns]
    
    print("Abriendo PowerPoint...")
    prs = Presentation(pptx_path)
    
    if not os.path.exists("temp_imgs"):
        os.mkdir("temp_imgs")

    # Contar registros a procesar
    valid_rows = []
    for idx, row in df.iterrows():
        razon = row.get('Razon_Social_Reclamante')
        if pd.isna(razon) or str(razon).strip() == "":
            continue
        if str(razon).strip() == "TOTAL GENERAL":
            break
        valid_rows.append((idx, row))
        
    for i, (idx, row) in enumerate(valid_rows):
        razon_social = row.get('Razon_Social_Reclamante')
        print(f"Procesando: {razon_social}")
        
        # Generar imagen (usando browser engine)
        styled = stylize_df(row, existing_cols)
        img_path = f"temp_imgs/img_{idx}.png"
        
        # Eliminamos table_conversion="matplotlib" para usar Chrome/Playwright y respetar CSS
        dfi.export(styled, img_path, dpi=200)
        
        # Diapositiva
        if i == 0:
            slide = prs.slides[0]
        else:
            slide = duplicate_slide(prs, 0)
        
        # Modificar título sin perder el formato (fuente, tamaño, negrita)
        title_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame and "Resumen Ejecutivo" in shape.text:
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    # Sobrescribimos el texto en el primer fragmento para mantener el estilo
                    shape.text_frame.paragraphs[0].runs[0].text = f"Resumen Ejecutivo de Gestión de\n{razon_social}"
                    # Vaciamos los demás si los hay
                    for k in range(1, len(shape.text_frame.paragraphs[0].runs)):
                        shape.text_frame.paragraphs[0].runs[k].text = ""
                else:
                    shape.text = f"Resumen Ejecutivo de Gestión de\n{razon_social}"
                title_shape = shape
                break
        
        # Insertar imagen 3 enters abajo
        if title_shape:
            top = title_shape.top + title_shape.height + Inches(0.5)
            left = Inches(1)
        else:
            top = Inches(2)
            left = Inches(1)
            
        slide.shapes.add_picture(img_path, left, top)

    print("Guardando presentación...")
    prs.save(out_pptx_path)
    print("¡Listo! Archivo guardado como", out_pptx_path)
    
    shutil.rmtree("temp_imgs")

if __name__ == "__main__":
    main()
